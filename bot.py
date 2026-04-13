import os
import re
import json
import base64
import requests
import datetime
import io
from telegram import Update
from telegram.ext import ApplicationBuilder, CommandHandler, MessageHandler, filters, ContextTypes

BOT_TOKEN  = os.environ.get("BOT_TOKEN", "")
CLAUDE_KEY = os.environ.get("CLAUDE_API_KEY", "")
BOT_PW     = os.environ.get("BOT_PASSWORD", "hanwha1234")
FB_URL     = os.environ.get("FB_URL", "https://team-dashboard-c0d7b-default-rtdb.asia-southeast1.firebasedatabase.app")
FB_SECRET  = os.environ.get("FB_SECRET", "")

AUTHED = set()

# ── Firebase ──────────────────────────────────────────────────
def fb_params():
    return {"auth": FB_SECRET} if FB_SECRET else {}

def fb_read(path):
    try:
        res = requests.get(f"{FB_URL}/{path}.json", params=fb_params(), timeout=5)
        return res.json() or {}
    except:
        return {}

def fb_write(path, data):
    try:
        res = requests.patch(f"{FB_URL}/{path}.json", json=data, params=fb_params(), timeout=5)
        return res.status_code == 200
    except:
        return False

def fb_push(path, data):
    try:
        requests.post(f"{FB_URL}/{path}.json", json=data, params=fb_params(), timeout=5)
    except:
        pass

# ── Claude API ────────────────────────────────────────────────
def ask_claude(messages, system="", max_tokens=1500):
    if not CLAUDE_KEY:
        return "AI 기능 준비 중"
    try:
        body = {"model": "claude-haiku-4-5", "max_tokens": max_tokens, "messages": messages}
        if system:
            body["system"] = system
        res = requests.post(
            "https://api.anthropic.com/v1/messages",
            headers={"x-api-key": CLAUDE_KEY, "anthropic-version": "2023-06-01", "content-type": "application/json"},
            json=body, timeout=60,
        )
        data = res.json()
        if "content" not in data:
            return f"오류: {data.get('error',{}).get('message','API 오류')}"
        return data["content"][0]["text"]
    except Exception as e:
        return f"오류: {str(e)}"

# ── 파일 텍스트 추출 ──────────────────────────────────────────
def extract_text_from_file(file_bytes, mime_type, file_name):
    """파일에서 텍스트 추출 (Excel, Word, PDF, PPT)"""
    ext = file_name.lower().split('.')[-1] if '.' in file_name else ''

    # Excel
    if ext in ('xlsx', 'xls') or 'excel' in mime_type or 'spreadsheet' in mime_type:
        try:
            import openpyxl
            wb = openpyxl.load_workbook(io.BytesIO(file_bytes), data_only=True)
            rows = []
            for sheet in wb.worksheets[:3]:  # 최대 3개 시트
                rows.append(f"[시트: {sheet.title}]")
                for row in sheet.iter_rows(min_row=1, max_row=100, values_only=True):
                    if any(c is not None for c in row):
                        rows.append('\t'.join([str(c) if c is not None else '' for c in row]))
            return '\n'.join(rows[:200])
        except ImportError:
            return None
        except Exception as e:
            return f"Excel 읽기 오류: {e}"

    # Word
    if ext in ('docx', 'doc') or 'word' in mime_type or 'document' in mime_type:
        try:
            import docx
            doc = docx.Document(io.BytesIO(file_bytes))
            return '\n'.join([p.text for p in doc.paragraphs if p.text.strip()])[:5000]
        except ImportError:
            # python-docx 없으면 텍스트 직접 추출 시도
            try:
                text = file_bytes.decode('utf-8', errors='ignore')
                return text[:3000]
            except:
                return None
        except Exception as e:
            return f"Word 읽기 오류: {e}"

    # PPT
    if ext in ('pptx', 'ppt') or 'powerpoint' in mime_type or 'presentation' in mime_type:
        try:
            from pptx import Presentation
            prs = Presentation(io.BytesIO(file_bytes))
            texts = []
            for i, slide in enumerate(prs.slides[:20]):
                texts.append(f"[슬라이드 {i+1}]")
                for shape in slide.shapes:
                    if hasattr(shape, 'text') and shape.text.strip():
                        texts.append(shape.text.strip())
            return '\n'.join(texts)[:5000]
        except ImportError:
            return None
        except Exception as e:
            return f"PPT 읽기 오류: {e}"

    # PDF → Claude Vision으로 처리 (base64 반환)
    if ext == 'pdf' or 'pdf' in mime_type:
        return 'PDF_BASE64:' + base64.b64encode(file_bytes).decode()

    return None

# ── Claude로 문서 분석 ────────────────────────────────────────
def analyze_document_with_claude(content, file_name, caption=""):
    """문서 내용을 Claude가 분석해서 구조화된 데이터 추출"""

    system = """당신은 HWR(Hanwha Q CELLS USA) 미국 재생에너지 개발 사업 문서 분석 전문가입니다.

업로드된 문서를 분석하여 다음 중 해당하는 데이터를 JSON으로 추출하세요.

1. 재무 데이터 (P&L/BS/CF):
{"type":"financial","stmt":"pl|bs|cf","year":2026,"month":3,"data":{"rev":124.0,"ebit":-3.0}}

2. 인허가 데이터:
{"type":"permit","project":"프로젝트명","permits":[{"key":"environmental","status":"승인","date":"2026-03-15","note":"메모"}]}

3. 프로젝트 현황:
{"type":"project","project":"프로젝트명","data":{"stage":"NBO","note":"내용","lead":"담당자"}}

4. PPV/투자 데이터:
{"type":"ppv","project":"프로젝트명","data":{"rf":0.8,"stage":"Late"}}

5. 전략/보고:
{"type":"report","category":"weekly|monthly|strategy","title":"제목","date":"2026-03-15","summary":"요약"}

6. 복합 데이터 (여러 항목):
{"type":"multi","items":[...]}

규칙:
- 억원 단위는 1380으로 나눠서 $M으로 변환
- 확인할 수 없는 항목은 포함하지 말것
- JSON만 반환, 다른 텍스트 없음
- 해당 데이터가 없으면: {"type":"unknown","summary":"문서 요약"}
"""

    user_msg = f"파일명: {file_name}\n"
    if caption:
        user_msg += f"사용자 설명: {caption}\n"
    user_msg += f"\n문서 내용:\n{content[:6000]}"

    return ask_claude([{"role": "user", "content": user_msg}], system=system, max_tokens=2000)

def analyze_pdf_with_claude(pdf_b64, file_name, caption=""):
    """PDF를 Claude Vision으로 분석"""
    system = """HWR 문서 분석 전문가. PDF를 분석하여 JSON으로 데이터를 추출하세요.
재무/인허가/프로젝트현황/PPV/보고서 중 해당 데이터를 추출.
억원→$M 변환(÷1380). JSON만 반환."""

    user_content = [
        {"type": "document", "source": {"type": "base64", "media_type": "application/pdf", "data": pdf_b64}},
        {"type": "text", "text": f"파일: {file_name}\n{caption}\n\n위 문서에서 데이터를 추출해 JSON으로 반환하세요."}
    ]
    return ask_claude([{"role": "user", "content": user_content}], system=system, max_tokens=2000)

def analyze_image_with_claude(img_b64, media_type, file_name, caption=""):
    """이미지를 Claude Vision으로 분석"""
    system = """HWR 문서 분석 전문가. 이미지에서 데이터를 추출하여 JSON으로 반환하세요.
재무/인허가/프로젝트현황/PPV/보고서 중 해당 데이터를 추출.
억원→$M 변환(÷1380). JSON만 반환."""

    user_content = [
        {"type": "image", "source": {"type": "base64", "media_type": media_type, "data": img_b64}},
        {"type": "text", "text": f"파일: {file_name}\n{caption}\n\n위 이미지에서 데이터를 추출해 JSON으로 반환하세요."}
    ]
    return ask_claude([{"role": "user", "content": user_content}], system=system, max_tokens=2000)

# ── Firebase 저장 라우터 ──────────────────────────────────────
def save_parsed_data(parsed):
    """파싱된 데이터를 Firebase에 저장"""
    dtype = parsed.get("type", "unknown")
    now = datetime.datetime.now().isoformat()[:16]
    results = []

    if dtype == "financial":
        stmt  = parsed.get("stmt")
        year  = parsed.get("year", 2026)
        month = parsed.get("month")
        data  = parsed.get("data", {})
        if stmt and month and data:
            data["updated_at"] = now
            if fb_write(f"financial/{stmt}/{year}/{month}", data):
                results.append(f"재무({stmt.upper()}) {year}년 {month}월 저장")

    elif dtype == "permit":
        project = parsed.get("project", "")
        permits = parsed.get("permits", [])
        for p in permits:
            key = p.get("key", "")
            if key:
                if fb_write(f"permits/{project}/{key}", {**p, "updated_at": now}):
                    results.append(f"인허가 [{project}] {key} 저장")

    elif dtype == "project":
        project = parsed.get("project", "")
        data    = parsed.get("data", {})
        if project and data:
            data["updated_at"] = now
            if fb_write(f"projects/{project}", data):
                results.append(f"프로젝트 [{project}] 현황 저장")

    elif dtype == "ppv":
        project = parsed.get("project", "")
        data    = parsed.get("data", {})
        if project and data:
            data["updated_at"] = now
            safe_name = project.replace("/", "_").replace(".", "_")
            if fb_write(f"ppv/overrides/{safe_name}", data):
                results.append(f"PPV [{project}] 저장")

    elif dtype == "report":
        category = parsed.get("category", "misc")
        fb_push(f"reports/{category}", {
            "title":   parsed.get("title", ""),
            "date":    parsed.get("date", now[:10]),
            "summary": parsed.get("summary", ""),
            "updated_at": now
        })
        results.append(f"보고서 [{category}] 저장")

    elif dtype == "multi":
        for item in parsed.get("items", []):
            sub = save_parsed_data(item)
            results.extend(sub)

    return results

# ── 자연어 재무 파싱 ──────────────────────────────────────────
PARSE_SYSTEM = """재무 데이터 파싱 전문가. 자연어에서 수치를 추출해 JSON으로만 반환.
{"stmt":"pl|bs|cf","year":2026,"month":3,"data":{"rev":124.0,"ebit":-3.0,"ni":-8.0}}
항목명 매핑: 매출→rev, 영업이익/EBIT→ebit, 순이익/당기순이익→ni, EBITDA→ebitda,
자산→ta, 부채→liab, 자본→equity, 현금→cash,
영업CF→ocf, 투자CF→icf, 재무CF→ffcf, 기말현금→end_cash
억원이면 1380으로 나눠서 $M으로. JSON만 반환."""

def parse_financial_text(text):
    result = ask_claude([{"role": "user", "content": text}], system=PARSE_SYSTEM, max_tokens=500)
    try:
        return json.loads(result.replace("```json","").replace("```","").strip())
    except:
        return None

def detect_intent(text):
    t = text.lower()
    fin_kw = ["업데이트","입력","저장","올려","넣어","실적","수치","매출","영업이익",
              "순이익","ebitda","자산","부채","자본","현금흐름","capex","월별","분기"]
    if any(k in t for k in fin_kw) and any(c.isdigit() for c in text):
        return "finance_update"
    return "question"

# ── HWR 컨텍스트 ──────────────────────────────────────────────
def get_hwr_context():
    ppv = fb_read("ppv/summary")
    fin = fb_read("financial")
    ctx = """답변 규칙: 마크다운 금지, 이모지 최소화, 보고서 형태로 간결하게.
당신은 HWR(Hanwha Q CELLS USA/HEUH) 미국 개발 사업 담당자입니다.

[포트폴리오] 94개 프로젝트, PV 13,761 MWac, ESS 10,871 MW
['26년 매각 $142M] Boulder Solar 3(H/$40M) Bonanza Peak(M/$50M) Oberon II~IV(M/$20M) Taormina/Lavender(M/$20M) Gibson(L/$12M)
[운영자산 유동화] 단기매각 $150~182M 현금유입, 차입금 $275M 제거
[Atlas Milestone] 달성률 25%, D-Day 임박
"""
    if ppv and ppv.get("totalRisked"):
        by = ppv.get("byStage", {})
        ctx += f"[PPV] 총 ${ppv.get('totalRisked')}M (Late ${by.get('Late',0)}M / Mid ${by.get('Mid',0)}M / Early ${by.get('Early',0)}M)\n"
    if fin:
        ctx += "\n[재무 데이터]\n"
        for stmt in ["pl","bs","cf"]:
            d = fin.get(stmt, {})
            if d:
                for y in sorted(d.keys(), reverse=True)[:1]:
                    for m in sorted(d[y].keys(), key=lambda x: int(x), reverse=True)[:3]:
                        items = {k:v for k,v in d[y][m].items() if k != "updated_at"}
                        ctx += f"  {stmt.upper()} {y}년 {m}월: {items}\n"
    return ctx

# ── 인증 ──────────────────────────────────────────────────────
def is_authed(uid): return uid in AUTHED

# ── 메뉴 ──────────────────────────────────────────────────────
async def show_menu(update):
    await update.message.reply_text(
        "HWR Dashboard Bot\n\n"
        "조회\n"
        "  /status    매각현황\n"
        "  /atlas     Atlas Milestone\n"
        "  /ppa       PPA 진척\n"
        "  /liquidity 운영자산 유동화\n"
        "  /strategy  전략 액션\n"
        "  /ppv       PPV 현황\n"
        "  /financial 재무 현황\n\n"
        "데이터 업로드\n"
        "  자연어: '26년 3월 매출 124, 영업이익 -3'\n"
        "  파일첨부: Excel, Word, PDF, PPT, 이미지\n"
        "  → Claude가 자동 파싱 후 Firebase 저장\n\n"
        "자유롭게 질문하셔도 됩니다."
    )

# ── 명령어 핸들러 ─────────────────────────────────────────────
async def start(update: Update, ctx: ContextTypes.DEFAULT_TYPE):
    if is_authed(update.effective_user.id): await show_menu(update)
    else: await update.message.reply_text("HWR Dashboard Bot입니다.\n비밀번호를 입력하세요.")

async def ppv_cmd(update: Update, ctx: ContextTypes.DEFAULT_TYPE):
    if not is_authed(update.effective_user.id):
        await update.message.reply_text("비밀번호를 먼저 입력하세요."); return
    data = fb_read("ppv")
    summary = data.get("summary", {})
    if not summary:
        await update.message.reply_text("PPV 데이터 없음\n대시보드에서 스냅샷을 찍어주세요."); return
    by = summary.get("byStage", {})
    msg = (f"PPV 현황 ({summary.get('updatedAt','?')[:10]})\n\n"
           f"총 Risked  ${summary.get('totalRisked','?')}M\n"
           f"  Late     ${by.get('Late',0):.1f}M\n"
           f"  Mid      ${by.get('Mid',0):.1f}M\n"
           f"  Early    ${by.get('Early',0):.1f}M")
    await update.message.reply_text(msg)

async def financial_cmd(update: Update, ctx: ContextTypes.DEFAULT_TYPE):
    if not is_authed(update.effective_user.id):
        await update.message.reply_text("비밀번호를 먼저 입력하세요."); return
    data = fb_read("financial")
    if not data:
        await update.message.reply_text("재무 데이터 없음\n자연어나 파일로 입력해주세요."); return
    msg = "재무 현황\n\n"
    for stmt, label in {"pl":"P&L","bs":"B/S","cf":"C/F"}.items():
        d = data.get(stmt, {})
        if not d: continue
        msg += f"{label}\n"
        for y in sorted(d.keys(), reverse=True)[:1]:
            for m in sorted(d[y].keys(), key=lambda x: int(x), reverse=True)[:3]:
                items = "  ".join([f"{k} {v}" for k,v in d[y][m].items() if k != "updated_at"])
                msg += f"  {y}년 {m}월  {items}\n"
        msg += "\n"
    await update.message.reply_text(msg)

async def status(update: Update, ctx: ContextTypes.DEFAULT_TYPE):
    if not is_authed(update.effective_user.id):
        await update.message.reply_text("비밀번호를 먼저 입력하세요."); return
    await update.message.reply_text(
        "'26년 매각현황  총 $142M\n\n"
        "H확도  Boulder Solar 3  $40M  NBO\n"
        "M확도  Bonanza Peak $50M  Oberon II~IV $20M  Taormina $10M  Lavender $10M\n"
        "L확도  Gibson $12M"
    )

async def atlas_cmd(update: Update, ctx: ContextTypes.DEFAULT_TYPE):
    if not is_authed(update.effective_user.id):
        await update.message.reply_text("비밀번호를 먼저 입력하세요."); return
    await update.message.reply_text(
        "Atlas North 1st Milestone  달성률 25%\n\n"
        "완료    2.12(f) CAP License\n"
        "진행    2.12(a) 2.12(c)\n"
        "D-Day   D+10 초과"
    )

async def liquidity_cmd(update: Update, ctx: ContextTypes.DEFAULT_TYPE):
    if not is_authed(update.effective_user.id):
        await update.message.reply_text("비밀번호를 먼저 입력하세요."); return
    await update.message.reply_text(
        "운영자산 유동화\n\n"
        "단기매각  TotalJV(Ob1A Rayos Ellis Skysol) HEUH(Laguna Astoria)\n"
        "현금유입 $150~182M  차입금제거 $275M  PL -$58~-90M\n\n"
        "중장기보유  Ho'Ohana  Oberon 1B  Imeson"
    )

async def strategy_cmd(update: Update, ctx: ContextTypes.DEFAULT_TYPE):
    if not is_authed(update.effective_user.id):
        await update.message.reply_text("비밀번호를 먼저 입력하세요."); return
    await update.message.reply_text(
        "전략 액션 아이템  이행률 14%\n\n"
        "1. 선제적 매각 프로세스   진행 중\n"
        "2. EPC Framework         완료\n"
        "3. Value-up 방안         진행 중\n"
        "4. Legacy 자산 전략      진행 중\n"
        "5. BESS 성장 전략        진행 중\n"
        "6. ISO별 Local GR        진행 중"
    )

async def ppa_cmd(update: Update, ctx: ContextTypes.DEFAULT_TYPE):
    if not is_authed(update.effective_user.id):
        await update.message.reply_text("비밀번호를 먼저 입력하세요."); return
    await update.message.reply_text(
        "PPA 현황  18개 프로젝트\n\n"
        "BL   Atlas 15\n"
        "RFP  Harlem River  Taormina  Lavender  Gibson\n"
        "미착수  13개  계약완료 0건"
    )

# ── 파일/이미지 핸들러 ────────────────────────────────────────
async def handle_file(update: Update, ctx: ContextTypes.DEFAULT_TYPE):
    if not is_authed(update.effective_user.id):
        await update.message.reply_text("비밀번호를 먼저 입력하세요."); return

    caption = update.message.caption or ""
    await update.message.reply_text("파일 분석 중... 잠시만 기다려 주세요.")

    try:
        # 이미지 처리
        if update.message.photo:
            photo = update.message.photo[-1]
            file = await ctx.bot.get_file(photo.file_id)
            file_bytes = await file.download_as_bytearray()
            img_b64 = base64.b64encode(file_bytes).decode()
            result = analyze_image_with_claude(img_b64, "image/jpeg", "image.jpg", caption)
            await process_result(update, ctx, result, "image.jpg")
            return

        if not update.message.document:
            await update.message.reply_text("파일을 첨부해주세요."); return

        doc = update.message.document
        file_name = doc.file_name or "unknown"
        mime_type = doc.mime_type or ""

        # 파일 크기 체크 (20MB 제한)
        if doc.file_size > 20 * 1024 * 1024:
            await update.message.reply_text("파일이 너무 큽니다. 20MB 이하만 가능합니다."); return

        file = await ctx.bot.get_file(doc.file_id)
        file_bytes = await file.download_as_bytearray()

        # 이미지 파일
        if mime_type.startswith("image/"):
            img_b64 = base64.b64encode(file_bytes).decode()
            result = analyze_image_with_claude(img_b64, mime_type, file_name, caption)
            await process_result(update, ctx, result, file_name)
            return

        # 텍스트 추출 시도
        extracted = extract_text_from_file(bytes(file_bytes), mime_type, file_name)

        if extracted is None:
            await update.message.reply_text(
                f"'{file_name}' 파일 형식을 지원하지 않습니다.\n"
                "지원: Excel(.xlsx), Word(.docx), PDF(.pdf), PPT(.pptx), 이미지"
            )
            return

        # PDF → Claude Vision
        if isinstance(extracted, str) and extracted.startswith("PDF_BASE64:"):
            pdf_b64 = extracted[11:]
            result = analyze_pdf_with_claude(pdf_b64, file_name, caption)
        else:
            # 텍스트 → Claude 분석
            result = analyze_document_with_claude(extracted, file_name, caption)

        await process_result(update, ctx, result, file_name)

    except Exception as e:
        await update.message.reply_text(f"파일 처리 중 오류가 발생했습니다.\n{str(e)}")

async def process_result(update, ctx, result, file_name):
    """Claude 분석 결과 처리"""
    try:
        clean = result.replace("```json","").replace("```","").strip()
        parsed = json.loads(clean)
    except:
        # JSON 파싱 실패 → 텍스트로 응답
        await update.message.reply_text(
            f"{file_name} 분석 완료\n\n{result[:1000]}"
        )
        return

    dtype = parsed.get("type", "unknown")

    if dtype == "unknown":
        await update.message.reply_text(
            f"{file_name} 분석 완료\n\n{parsed.get('summary', '데이터를 찾을 수 없습니다.')}"
        )
        return

    # 저장 전 확인
    preview = format_preview(parsed)
    await update.message.reply_text(
        f"{file_name} 분석 완료\n\n{preview}\n\n저장하려면 '확인', 취소하려면 '취소'"
    )
    ctx.user_data["pending_parsed"] = parsed

def format_preview(parsed):
    """파싱 결과 미리보기"""
    dtype = parsed.get("type")
    if dtype == "financial":
        stmt = parsed.get("stmt","?").upper()
        year = parsed.get("year","?")
        month = parsed.get("month","?")
        items = "\n".join([f"  {k}: {v}" for k,v in parsed.get("data",{}).items()])
        return f"재무 데이터 ({stmt} {year}년 {month}월)\n{items}"
    elif dtype == "permit":
        proj = parsed.get("project","?")
        permits = parsed.get("permits",[])
        lines = [f"인허가 [{proj}]"]
        for p in permits:
            lines.append(f"  {p.get('key','?')}: {p.get('status','?')} {p.get('date','')}")
        return "\n".join(lines)
    elif dtype == "project":
        proj = parsed.get("project","?")
        data = parsed.get("data",{})
        items = "\n".join([f"  {k}: {v}" for k,v in data.items()])
        return f"프로젝트 [{proj}]\n{items}"
    elif dtype == "ppv":
        proj = parsed.get("project","?")
        data = parsed.get("data",{})
        return f"PPV [{proj}]\n" + "\n".join([f"  {k}: {v}" for k,v in data.items()])
    elif dtype == "report":
        return f"보고서 [{parsed.get('category','?')}]\n  제목: {parsed.get('title','?')}\n  날짜: {parsed.get('date','?')}"
    elif dtype == "multi":
        items = parsed.get("items",[])
        return f"복합 데이터 {len(items)}건\n" + "\n".join([format_preview(i) for i in items[:3]])
    return str(parsed)[:500]

# ── 메시지 핸들러 ─────────────────────────────────────────────
async def handle_message(update: Update, ctx: ContextTypes.DEFAULT_TYPE):
    uid  = update.effective_user.id
    text = update.message.text.strip()

    if not is_authed(uid):
        if text == BOT_PW:
            AUTHED.add(uid)
            await update.message.reply_text("인증되었습니다.")
            await show_menu(update)
        else:
            await update.message.reply_text("비밀번호가 틀렸습니다. /start")
        return

    # 확인/취소
    if text in ("확인","저장","yes","ㅇ","ㅇㅇ"):
        pending = ctx.user_data.get("pending_parsed") or ctx.user_data.get("pending_fin")
        if pending:
            results = save_parsed_data(pending)
            if results:
                await update.message.reply_text("저장 완료\n" + "\n".join(["  " + r for r in results]) + "\n\n대시보드에서 확인하세요.")
            else:
                await update.message.reply_text("저장 실패. 다시 시도해주세요.")
            ctx.user_data.pop("pending_parsed", None)
            ctx.user_data.pop("pending_fin", None)
        else:
            await update.message.reply_text("저장할 데이터가 없습니다.")
        return

    if text in ("취소","no","ㄴ"):
        ctx.user_data.pop("pending_parsed", None)
        ctx.user_data.pop("pending_fin", None)
        await update.message.reply_text("취소했습니다.")
        return

    # 재무 자연어 입력
    if detect_intent(text) == "finance_update":
        await update.message.reply_text("파싱 중...")
        parsed_fin = parse_financial_text(text)
        if parsed_fin and all(k in parsed_fin for k in ["stmt","year","month","data"]):
            parsed_fin["type"] = "financial"
            preview = format_preview(parsed_fin)
            await update.message.reply_text(f"아래 내용으로 저장할까요?\n\n{preview}\n\n저장하려면 '확인', 취소하려면 '취소'")
            ctx.user_data["pending_parsed"] = parsed_fin
        else:
            await update.message.reply_text(
                "데이터 인식 실패\n\n예시:\n"
                "'26년 3월 PL: 매출 124, 영업이익 -3, 순이익 -8\n"
                "'2026년 3월 BS: 자산 450, 부채 320, 자본 130"
            )
        return

    # 일반 질문
    await update.message.reply_text("분석 중...")
    answer = ask_claude([{"role":"user","content":text}], system=get_hwr_context(), max_tokens=800)
    await update.message.reply_text(answer)

# ── 실행 ──────────────────────────────────────────────────────
if __name__ == "__main__":
    app = ApplicationBuilder().token(BOT_TOKEN).build()
    app.add_handler(CommandHandler("start",     start))
    app.add_handler(CommandHandler("status",    status))
    app.add_handler(CommandHandler("atlas",     atlas_cmd))
    app.add_handler(CommandHandler("liquidity", liquidity_cmd))
    app.add_handler(CommandHandler("strategy",  strategy_cmd))
    app.add_handler(CommandHandler("ppa",       ppa_cmd))
    app.add_handler(CommandHandler("ppv",       ppv_cmd))
    app.add_handler(CommandHandler("financial", financial_cmd))
    app.add_handler(MessageHandler(filters.PHOTO, handle_file))
    app.add_handler(MessageHandler(filters.Document.ALL, handle_file))
    app.add_handler(MessageHandler(filters.TEXT & ~filters.COMMAND, handle_message))
    print("HWR Bot 시작 (Excel/Word/PDF/PPT/이미지 파싱)")
    app.run_polling()
